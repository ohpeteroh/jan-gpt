import os
import tempfile
from pathlib import Path
import streamlit as st
import pandas as pd
import subprocess
from pptx import Presentation
import pyxlsb
import pytesseract
from PIL import Image
from duckduckgo_search import DDGS

from langchain.docstore.document import Document
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings, ChatOpenAI

# 🔐 API Key
OPENAI_API_KEY = st.secrets["OPENAI_API_KEY"]

embedding = OpenAIEmbeddings(openai_api_key=OPENAI_API_KEY)
llm = ChatOpenAI(model="gpt-4", temperature=0.2, openai_api_key=OPENAI_API_KEY)
splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
DB_PATH = "faiss_index"

def load_and_split_file(tmp_path, suffix):
    docs = []
    if suffix == ".txt":
        with open(tmp_path, encoding="utf-8") as f:
            docs = [Document(page_content=f.read())]
    elif suffix == ".pdf":
        from langchain_community.document_loaders import PyPDFLoader
        docs = PyPDFLoader(tmp_path).load()
    elif suffix == ".docx":
        from langchain_community.document_loaders import UnstructuredWordDocumentLoader
        docs = UnstructuredWordDocumentLoader(tmp_path).load()
    elif suffix == ".pptx":
        prs = Presentation(tmp_path)
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        docs = [Document(page_content=text)]
    elif suffix == ".hwp":
        result = subprocess.run(['hwp5txt', tmp_path], stdout=subprocess.PIPE, encoding='utf-8')
        docs = [Document(page_content=result.stdout)]
    elif suffix in [".xlsx", ".xlsm"]:
        df = pd.read_excel(tmp_path, engine='openpyxl')
        docs = [Document(page_content=df.to_string())]
    elif suffix == ".xlsb":
        with pyxlsb.open_workbook(tmp_path) as wb:
            sheet = wb.get_sheet(1)
            data = "\n".join(["\t".join([str(cell.v) for cell in row]) for row in sheet.rows()])
        docs = [Document(page_content=data)]
    elif suffix in [".csv"]:
        df = pd.read_csv(tmp_path)
        docs = [Document(page_content=df.to_string())]
    elif suffix in [".png", ".jpg", ".jpeg"]:
        text = pytesseract.image_to_string(Image.open(tmp_path), lang='eng+kor')
        docs = [Document(page_content=text)]

    if not docs:
        return False

    chunks = splitter.split_documents(docs)
    if not os.path.exists(DB_PATH):
        db = FAISS.from_documents(chunks, embedding)
    else:
        db = FAISS.load_local(DB_PATH, embedding)
        db.add_documents(chunks)
    db.save_local(DB_PATH)
    return True

st.set_page_config(page_title="Jan GPT", layout="wide")
st.title("📂 Jan GPT - 문서 + 이미지 + 웹 검색 기반 리서치 GPT")

uploaded_file = st.file_uploader("📤 파일 업로드 (.txt, .pdf, .docx, .pptx, .hwp, .xlsx, .xlsm, .xlsb, .csv, .png, .jpg)",
                                 type=["txt", "pdf", "docx", "pptx", "hwp", "xlsx", "xlsm", "xlsb", "csv", "png", "jpg", "jpeg"])

if uploaded_file:
    suffix = Path(uploaded_file.name).suffix.lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(uploaded_file.getvalue())
        tmp_path = tmp.name
    if load_and_split_file(tmp_path, suffix):
        st.success("✅ 파일이 자동으로 벡터화되어 학습되었습니다.")

query = st.text_input("❓ 질문을 입력하세요")
use_web = st.checkbox("🌐 웹 검색도 포함할까요? (DuckDuckGo 기반)", value=False)
search_mode = st.radio("검색 모드", ["일반 검색", "심층 리서치"], horizontal=True)

if query:
    try:
        if os.path.exists(DB_PATH):
            db = FAISS.load_local(DB_PATH, embedding)
            docs = db.similarity_search(query, k=5)
            doc_context = "\n\n".join([doc.page_content for doc in docs])
        else:
            doc_context = "(문서 없음)"

        # DuckDuckGo 검색
        web_results = ""
        if use_web:
            try:
                ddgs = DDGS()
                results = ddgs.text(query, max_results=5)
                web_results = "\n".join([r["body"] for r in results])
            except Exception as e:
                web_results = f"(웹 검색 실패: {e})"

        # 프롬프트 구성
        prompt = f"[문서 기반 정보]\n{doc_context}\n"
        if use_web:
            prompt += f"\n[웹 검색 정보]\n{web_results}\n"

        if search_mode == "심층 리서치":
            prompt += f"\n위 정보를 바탕으로 '{query}'에 대해 다음 항목을 포함한 심층 분석 보고서를 작성해주세요:\n1. 핵심 요약\n2. 주요 근거 및 배경 정보\n3. 전략적 시사점 및 제언"
        else:
            prompt += f"\n위 정보를 바탕으로 '{query}'에 답변해 주세요."

        response = llm.invoke(prompt)
        st.markdown("### 💬 GPT 응답")
        st.write(response.content)

    except Exception as e:
        st.error(f"❌ 오류 발생: {e}")
